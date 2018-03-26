import pptcom
import pytest
from pptx import Presentation
import os


@pytest.fixture
def presentation(tmpdir):
    filename = os.path.join(tmpdir, 'test.pptx')

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"

    prs.save(filename)
    return filename


@pytest.mark.parametrize('export_format', pptcom.files.EXPORT_FORMATS)
def test_export(presentation, export_format, tmpdir):
    with pptcom.File(presentation) as pptfile:
        pptfile.export(export_format)
        pptfile.export(export_format, path=tmpdir)
